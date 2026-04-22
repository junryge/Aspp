"""
demos_v1/ppt_builder.py - MD → .pptx 결정론적 변환기 (LLM 없이)

사용자가 MD 파일/텍스트를 입력하면 python-pptx 로 직접 렌더링.
LLM 이 끼어들지 않으므로 코드 에러·토큰 제한 등 문제 없음.

MD 규칙:
- # 제목       → 제목 슬라이드 (부제는 다음 줄의 > 텍스트)
- ## 슬라이드  → 내용 슬라이드 시작
- - 불릿       → 본문 불릿 (들여쓰기로 level 조절)
- | 표 |       → 표 슬라이드 (자동 감지)
- ```lang\ncode\n``` → 코드 슬라이드 (모노스페이스)
- ![캡션](경로) → 이미지 슬라이드 (Phase 2)
- --- 또는 ## 재등장 → 슬라이드 분리
"""
import io
import os
import re
import time
import uuid
from typing import Any

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False


# ============================================================
# 한글 폰트 폴백 체인
# ============================================================
KOREAN_FONT_CHAIN = ["맑은 고딕", "Malgun Gothic", "NanumGothic",
                    "Noto Sans KR", "Pretendard", "Arial"]
CODE_FONT_CHAIN = ["Consolas", "D2Coding", "Courier New", "Monaco"]


# ============================================================
# MD Parser
# ============================================================
def parse_md_to_outline(md_text: str) -> dict:
    """MD 텍스트 → outline dict 변환.

    Returns:
        {
            "meta": {"title": str, "subtitle": str},
            "slides": [
                {"type": "title", "title": str, "subtitle": str},
                {"type": "content", "title": str, "bullets": [{"text", "level"}]},
                {"type": "table", "title": str, "headers": [...], "rows": [[...]]},
                {"type": "code", "title": str, "language": str, "code": str},
                ...
            ]
        }
    """
    lines = md_text.replace("\r\n", "\n").split("\n")
    meta = {"title": "", "subtitle": ""}
    slides: list[dict] = []
    current: dict | None = None  # 현재 빌드 중인 슬라이드
    in_code_block = False
    code_lang = ""
    code_buffer: list[str] = []
    table_buffer: list[str] = []
    in_table = False

    def flush_table():
        """누적된 table_buffer 를 table 슬라이드로 변환 + current 로 승격."""
        nonlocal current, table_buffer, in_table
        if not table_buffer:
            return
        # header | separator | rows
        rows_raw = [l for l in table_buffer if l.strip().startswith("|")]
        if len(rows_raw) < 2:
            table_buffer = []
            in_table = False
            return

        def _split_row(row: str) -> list[str]:
            row = row.strip().strip("|")
            return [c.strip() for c in row.split("|")]

        headers = _split_row(rows_raw[0])
        # rows_raw[1] 은 구분자 (|---|---|) → skip
        body_rows = [_split_row(r) for r in rows_raw[2:] if r.strip()]
        title = current.get("title", "") if current else ""
        table_slide = {"type": "table", "title": title,
                       "headers": headers, "rows": body_rows}
        # 기존 current 가 content slide 였고 bullet 없으면 table 로 대체
        if current and current.get("type") == "content" and not current.get("bullets"):
            slides[-1] = table_slide
            current = table_slide
        else:
            slides.append(table_slide)
            current = table_slide
        table_buffer = []
        in_table = False

    for raw in lines:
        line = raw.rstrip()

        # 코드 블록 내부
        if in_code_block:
            if line.startswith("```"):
                # 코드 블록 종료 → 슬라이드로 확정
                code_slide = {
                    "type": "code",
                    "title": current.get("title", "") if current else "코드",
                    "language": code_lang,
                    "code": "\n".join(code_buffer),
                }
                if current and current.get("type") == "content" and not current.get("bullets"):
                    slides[-1] = code_slide
                else:
                    slides.append(code_slide)
                current = code_slide
                in_code_block = False
                code_lang = ""
                code_buffer = []
            else:
                code_buffer.append(raw)  # raw 쓰기 (들여쓰기 유지)
            continue

        # 테이블 버퍼링
        if line.strip().startswith("|"):
            in_table = True
            table_buffer.append(line)
            continue
        elif in_table:
            flush_table()

        # 코드 블록 시작
        if line.startswith("```"):
            in_code_block = True
            code_lang = line[3:].strip()
            code_buffer = []
            continue

        # --- 구분자
        if line.strip() in ("---", "***"):
            current = None
            continue

        # # 제목 (문서 메타)
        m1 = re.match(r"^#\s+(.+)$", line)
        if m1:
            title_text = m1.group(1).strip()
            if not slides:
                meta["title"] = title_text
                slide = {"type": "title", "title": title_text, "subtitle": ""}
                slides.append(slide)
                current = slide
            else:
                slide = {"type": "title", "title": title_text, "subtitle": ""}
                slides.append(slide)
                current = slide
            continue

        # ## 슬라이드 시작
        m2 = re.match(r"^##\s+(.+)$", line)
        if m2:
            slide = {"type": "content", "title": m2.group(1).strip(), "bullets": []}
            slides.append(slide)
            current = slide
            continue

        # ### ~ ###### 는 하위 제목으로 쓰거나 불릿 level 로
        m3 = re.match(r"^(#{3,6})\s+(.+)$", line)
        if m3:
            # 현재 슬라이드에 sub-heading 불릿으로 추가
            if current and current.get("type") == "content":
                level = len(m3.group(1)) - 2  # ### → 1, #### → 2
                current["bullets"].append({"text": m3.group(2).strip(), "level": level, "bold": True})
            continue

        # > 인용 (title 슬라이드의 subtitle 또는 본문 인용)
        mq = re.match(r"^>\s?(.*)$", line)
        if mq and current:
            if current.get("type") == "title" and not current.get("subtitle"):
                current["subtitle"] = mq.group(1).strip()
            elif current.get("type") == "content":
                current["bullets"].append({"text": "“" + mq.group(1).strip() + "”",
                                           "level": 0, "italic": True})
            continue

        # - 불릿 / 1. 번호 매김
        mb = re.match(r"^(\s*)([-*+]|\d+\.)\s+(.+)$", line)
        if mb:
            indent = len(mb.group(1))
            level = indent // 2   # 2 space = 1 level
            text = mb.group(3).strip()
            if not current or current.get("type") != "content":
                # 상위 슬라이드 없으면 "내용" 슬라이드 자동 생성
                slide = {"type": "content", "title": "내용", "bullets": []}
                slides.append(slide)
                current = slide
            current.setdefault("bullets", []).append({"text": text, "level": level})
            continue

        # 빈 줄은 무시
        if not line.strip():
            continue

        # 그 외 일반 텍스트 → 현재 슬라이드 본문에 추가
        if current and current.get("type") == "content":
            current.setdefault("bullets", []).append({"text": line.strip(), "level": 0})
        elif current and current.get("type") == "title" and not current.get("subtitle"):
            current["subtitle"] = line.strip()

    # 루프 종료 — 남은 테이블 flush
    if in_table:
        flush_table()

    # 슬라이드 하나도 없으면 안내 슬라이드
    if not slides:
        slides.append({"type": "title", "title": "빈 문서",
                       "subtitle": "MD 내용이 비어 있습니다."})

    return {"meta": meta, "slides": slides}


# ============================================================
# Renderer (python-pptx)
# ============================================================
def _apply_korean_font(run, bold: bool = False, italic: bool = False,
                      size_pt: int = 18, is_code: bool = False):
    """run 에 한글 친화 폰트 + 속성 적용."""
    font = run.font
    font.name = (CODE_FONT_CHAIN[0] if is_code else KOREAN_FONT_CHAIN[0])
    font.size = Pt(size_pt)
    font.bold = bold
    font.italic = italic


def _add_title_slide(prs, slide_data: dict):
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    if title:
        title.text = slide_data.get("title", "")
        for p in title.text_frame.paragraphs:
            for r in p.runs:
                _apply_korean_font(r, bold=True, size_pt=44)
    # 부제 placeholder (idx 1)
    subtitle_text = slide_data.get("subtitle", "")
    if subtitle_text and len(slide.placeholders) > 1:
        sub = slide.placeholders[1]
        sub.text = subtitle_text
        for p in sub.text_frame.paragraphs:
            for r in p.runs:
                _apply_korean_font(r, size_pt=22)
    return slide


def _add_content_slide(prs, slide_data: dict):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    if title:
        title.text = slide_data.get("title", "")
        for p in title.text_frame.paragraphs:
            for r in p.runs:
                _apply_korean_font(r, bold=True, size_pt=32)
    # 본문 placeholder
    body = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0:  # title 제외
            body = ph
            break
    if body is None:
        return slide
    tf = body.text_frame
    tf.word_wrap = True
    bullets = slide_data.get("bullets", [])
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = ""  # 기본 텍스트 지우기
        else:
            p = tf.add_paragraph()
        p.level = min(b.get("level", 0), 4)
        run = p.add_run()
        run.text = b.get("text", "")
        _apply_korean_font(run, bold=b.get("bold", False),
                          italic=b.get("italic", False), size_pt=18)
    return slide


def _add_table_slide(prs, slide_data: dict):
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    if title:
        title.text = slide_data.get("title", "") or "표"
        for p in title.text_frame.paragraphs:
            for r in p.runs:
                _apply_korean_font(r, bold=True, size_pt=32)
    headers = slide_data.get("headers", [])
    rows = slide_data.get("rows", [])
    if not headers:
        return slide
    n_cols = len(headers)
    n_rows = len(rows) + 1
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9.0)
    height = Inches(0.4 * n_rows + 0.2)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table
    # 헤더
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = str(h)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                _apply_korean_font(r, bold=True, size_pt=14)
    # 본문
    for i, row in enumerate(rows, start=1):
        for j in range(n_cols):
            cell = tbl.cell(i, j)
            cell.text = str(row[j]) if j < len(row) else ""
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    _apply_korean_font(r, size_pt=12)
    return slide


def _add_code_slide(prs, slide_data: dict):
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    if title:
        lang = slide_data.get("language", "")
        title.text = (slide_data.get("title", "") or "코드") + (f" ({lang})" if lang else "")
        for p in title.text_frame.paragraphs:
            for r in p.runs:
                _apply_korean_font(r, bold=True, size_pt=28)
    # 코드 박스
    from pptx.enum.shapes import MSO_SHAPE
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(5.5)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
    box.line.color.rgb = RGBColor(0x44, 0x44, 0x44)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.1)
    code = slide_data.get("code", "")
    lines = code.split("\n") or [""]
    for i, ln in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = ""
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = ln
        _apply_korean_font(run, size_pt=11, is_code=True)
        run.font.color.rgb = RGBColor(0xDC, 0xDC, 0xDC)
    return slide


def _make_presentation():
    """python-pptx Presentation 생성. 내장 default.pptx 가 누락된 환경 대응.

    python-pptx 설치가 깨져 default.pptx 가 없을 때 대비해:
      1. 우리 번들 템플릿(corporate.pptx 등) 시도
      2. 모두 실패하면 빈 상태로 수동 생성 시도
      3. 그래도 안 되면 원본 에러 그대로 raise
    """
    _HERE = os.path.dirname(os.path.abspath(__file__))
    _TEMPLATE_CANDIDATES = [
        os.path.join(_HERE, "ppt_templates", "corporate.pptx"),
        os.path.join(_HERE, "ppt_templates", "minimal.pptx"),
        os.path.join(_HERE, "ppt_templates", "academic.pptx"),
        os.path.join(_HERE, "ppt_templates", "creative.pptx"),
        os.path.join(_HERE, "ppt_templates", "dark.pptx"),
    ]
    # 1) python-pptx 기본 시도
    try:
        return Presentation()
    except Exception as e0:
        last_err = e0
        print(f"[ppt_builder] Presentation() 기본값 실패: {e0}")
    # 2) 번들 템플릿 순서대로 시도
    for tpl in _TEMPLATE_CANDIDATES:
        if not os.path.isfile(tpl):
            continue
        try:
            prs = Presentation(tpl)
            # 템플릿의 기존 슬라이드 제거 (샘플)
            try:
                xml_slides = prs.slides._sldIdLst
                for sld in list(xml_slides):
                    xml_slides.remove(sld)
            except Exception:
                pass
            print(f"[ppt_builder] 템플릿 폴백 사용: {os.path.basename(tpl)}")
            return prs
        except Exception as et:
            last_err = et
            print(f"[ppt_builder] 템플릿 시도 실패 {os.path.basename(tpl)}: {et}")
    # 3) 모두 실패 → 명확한 안내 메시지
    raise RuntimeError(
        "python-pptx 설치가 손상되어 default.pptx 가 없고 번들 템플릿도 로드 실패.\n"
        "해결: pip uninstall python-pptx -y && pip install python-pptx\n"
        f"원본 에러: {last_err}"
    )


def render_outline_to_pptx(outline: dict) -> bytes:
    """outline dict → .pptx bytes 반환."""
    if not _PPTX_AVAILABLE:
        raise RuntimeError("python-pptx 가 설치되지 않았습니다. pip install python-pptx")
    prs = _make_presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    dispatch = {
        "title": _add_title_slide,
        "content": _add_content_slide,
        "table": _add_table_slide,
        "code": _add_code_slide,
    }
    for slide_data in outline.get("slides", []):
        fn = dispatch.get(slide_data.get("type", "content"), _add_content_slide)
        try:
            fn(prs, slide_data)
        except Exception as e:
            # 에러 슬라이드로 대체 (전체 실패 방지)
            err_slide = prs.slides.add_slide(prs.slide_layouts[5])
            if err_slide.shapes.title:
                err_slide.shapes.title.text = f"[렌더링 실패] {slide_data.get('type', '?')}"
            print(f"[ppt_builder] slide render error: {e}")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ============================================================
# 파일 저장 + ID 관리
# ============================================================
PPT_CACHE_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                             "..", "uploads", "ppt_cache")
os.makedirs(PPT_CACHE_DIR, exist_ok=True)


def save_pptx(pptx_bytes: bytes, hint: str = "presentation") -> dict:
    """.pptx bytes 를 임시 파일로 저장하고 id/경로 반환."""
    ppt_id = f"{int(time.time())}_{uuid.uuid4().hex[:8]}"
    safe_hint = re.sub(r"[^\w가-힣]+", "_", hint)[:40] or "presentation"
    fname = f"{safe_hint}_{ppt_id}.pptx"
    path = os.path.join(PPT_CACHE_DIR, fname)
    with open(path, "wb") as f:
        f.write(pptx_bytes)
    return {"id": ppt_id, "filename": fname, "path": path, "size": len(pptx_bytes)}


def get_pptx_path(ppt_id: str) -> str | None:
    """ID 로 .pptx 파일 경로 조회."""
    if not ppt_id or "/" in ppt_id or "\\" in ppt_id or ".." in ppt_id:
        return None
    for fname in os.listdir(PPT_CACHE_DIR):
        if ppt_id in fname and fname.endswith(".pptx"):
            return os.path.join(PPT_CACHE_DIR, fname)
    return None


# ============================================================
# 편의 함수: md_text 한 번에 .pptx 로
# ============================================================
def md_to_pptx_file(md_text: str, title_hint: str = "") -> dict:
    """MD 텍스트 → outline → .pptx 파일 저장 → id 반환."""
    outline = parse_md_to_outline(md_text)
    pptx_bytes = render_outline_to_pptx(outline)
    title = title_hint or outline.get("meta", {}).get("title", "presentation")
    info = save_pptx(pptx_bytes, hint=title)
    info["outline"] = outline
    info["slide_count"] = len(outline.get("slides", []))
    return info
